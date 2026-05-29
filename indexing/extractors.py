"""Content extraction for indexing."""

from __future__ import annotations

import io
import logging
import os
from typing import Optional

import filetype
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from PIL import Image

from security import limits

logger = logging.getLogger(__name__)

_TEXT_MIME_PREFIXES = ("text/",)
_DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_PDF_MIME = "application/pdf"


def _tesseract_available() -> bool:
    try:
        import pytesseract

        if os.environ.get("TESSERACT_CMD"):
            pytesseract.pytesseract.tesseract_cmd = os.environ["TESSERACT_CMD"]
        return True
    except Exception:
        return False


def _extract_pdf(file_bytes: bytes) -> str:
    text_parts: list[str] = []
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    for page in doc:
        text_parts.append(page.get_text("text") or "")
    return "\n".join(text_parts)


def _extract_docx(file_bytes: bytes) -> str:
    doc = Document(io.BytesIO(file_bytes))
    return "\n".join(p.text for p in doc.paragraphs if p.text)


def _extract_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    lines: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text_frame.text or ""
                if text:
                    lines.append(text)
    return "\n".join(lines)


def _extract_text(file_bytes: bytes) -> str:
    try:
        return file_bytes.decode("utf-8", errors="ignore")
    except Exception:
        return ""


def _extract_image_ocr(file_bytes: bytes) -> str:
    if not _tesseract_available():
        logger.info("OCR skipped: pytesseract not available.")
        return ""
    import pytesseract

    try:
        img = Image.open(io.BytesIO(file_bytes))
        return pytesseract.image_to_string(img)
    except Exception:
        logger.warning("OCR failed for image content.")
        return ""


def detect_mime(file_bytes: bytes, filename: Optional[str] = None) -> str:
    guess = filetype.guess(file_bytes)
    if guess:
        return guess.mime
    if filename:
        ext = os.path.splitext(filename)[1].lower()
        if ext == ".pdf":
            return _PDF_MIME
        if ext == ".docx":
            return _DOCX_MIME
        if ext == ".pptx":
            return _PPTX_MIME
    return "application/octet-stream"


def extract_text(file_bytes: bytes, mime_type: str, filename: Optional[str] = None) -> str:
    if not file_bytes:
        return ""
    if limits.MAX_INDEX_BYTES and len(file_bytes) > limits.MAX_INDEX_BYTES:
        return ""

    mime = (mime_type or "").lower()
    try:
        if mime == _PDF_MIME:
            return _extract_pdf(file_bytes)
        if mime == _DOCX_MIME:
            return _extract_docx(file_bytes)
        if mime == _PPTX_MIME:
            return _extract_pptx(file_bytes)
        if mime.startswith(_TEXT_MIME_PREFIXES) or mime in {"application/json"}:
            return _extract_text(file_bytes)
        if mime.startswith("image/"):
            return _extract_image_ocr(file_bytes)
    except Exception:
        logger.warning("Text extraction failed for mime=%s", mime)

    # Try best-effort detection fallback
    detected = detect_mime(file_bytes, filename)
    if detected != mime:
        return extract_text(file_bytes, detected, filename)
    return ""
